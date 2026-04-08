import collections 
import collections.abc

# Fix for deprecated collections.abc module aliases in python 3.10+
try:
    collections.Iterator = collections.abc.Iterator
    collections.Mapping = collections.abc.Mapping
    collections.Sequence = collections.abc.Sequence
except AttributeError:
    pass

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def apply_dark_theme(slide):
    """Applies a dark cyber theme background with a cyan accent line to a slide."""
    # Add dark background shape
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(14, 18, 25)  # Very dark slate/blue
    bg.line.fill.background() # No border
    # Send background to back by placing it at the start of the shape tree
    slide.shapes.element.insert(0, bg.element)

    # Add Cyan top accent bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.1)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0, 240, 255) # Cyber Cyan
    accent.line.fill.background()


def format_title(title_shape, text, size=40, align=PP_ALIGN.LEFT):
    title_shape.text = text
    title_shape.text_frame.paragraphs[0].alignment = align
    for p in title_shape.text_frame.paragraphs:
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(0, 240, 255) # Cyan
        p.font.bold = True

def add_body_text(slide, points, left, top, width, height, font_size=20):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(points):
        p = tf.add_paragraph()
        if isinstance(point, tuple): # Header + Subtext
            p.text = point[0]
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(font_size + 2)
            
            p_sub = tf.add_paragraph()
            p_sub.text = "   " + point[1]
            p_sub.font.color.rgb = RGBColor(180, 180, 190)
            p_sub.font.size = Pt(font_size - 2)
            p_sub.space_after = Pt(20)
        else:
            p.text = "• " + point
            p.font.color.rgb = RGBColor(240, 240, 245)
            p.font.size = Pt(font_size)
            p.space_after = Pt(14)
    return tf

prs = Presentation()
# Set presentation size to widescreen 16:9
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6] # Blank layout removes default placeholders

# --- SLIDE 1: Title Slide ---
slide1 = prs.slides.add_slide(blank_layout)
apply_dark_theme(slide1)
# Center elements
tBox = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(1))
tf = tBox.text_frame
p = tf.add_paragraph()
p.text = "CCAF Firewall Engine"
p.font.bold = True
p.font.size = Pt(64)
p.font.color.rgb = RGBColor(0, 240, 255)
p.alignment = PP_ALIGN.CENTER

subBox = slide1.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11.3), Inches(1))
sub_tf = subBox.text_frame
sub_p = sub_tf.add_paragraph()
sub_p.text = "Next-Generation Stateful Traffic Analysis & Edge Security"
sub_p.font.size = Pt(28)
sub_p.font.color.rgb = RGBColor(200, 200, 210)
sub_p.alignment = PP_ALIGN.CENTER


# --- SLIDE 2: The Problem ---
slide2 = prs.slides.add_slide(blank_layout)
apply_dark_theme(slide2)
title2 = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
format_title(title2, "The Market Problem")
add_body_text(slide2, [
    ("Legacy Firewalls are Blind", "Traditional network defenses rely on static rules. They lack contextual awareness of real-time connection states."),
    ("High Latency in Proxies", "Modern deep-packet inspection often requires funneling traffic through sluggish proprietary proxies."),
    ("Poor Observability", "Operators are forced to dig through archaic CLI interfaces or dense logs rather than reacting to visual threat intelligence in real-time.")
], Inches(1), Inches(1.8), Inches(11), Inches(5))


# --- SLIDE 3: The Solution (CCAF) ---
slide3 = prs.slides.add_slide(blank_layout)
apply_dark_theme(slide3)
title3 = slide3.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
format_title(title3, "The Solution: Fast, Stateful, Observable")
add_body_text(slide3, [
    ("Zero-Latency Execution", "CCAF operates locally directly on the OS kernel's socket layer, making pass/drop decisions in milliseconds."),
    ("Stateful Track Analytics", "Tracks full TCP lifecycles (SYN, ACK, FIN) instead of just blocking port protocols."),
    ("Unprecedented Observability", "A highly optimized dashboard visualizes real-time load, active tracking, and dropped events instantly.")
], Inches(1), Inches(1.8), Inches(11), Inches(5))


# --- SLIDE 4: Core Technology Stack ---
slide4 = prs.slides.add_slide(blank_layout)
apply_dark_theme(slide4)
title4 = slide4.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
format_title(title4, "Proprietary Architecture")
add_body_text(slide4, [
    ("Decoupled MVC Pattern", "The inspection engine is heavily segregated from the UI, ensuring that visual rendering never bottlenecks packet processing."),
    ("The Stateful Engine (Backend)", "Python & psutil power a heavily threaded, in-memory state table capturing live system sockets."),
    ("The Command Matrix (Frontend)", "A Flask-driven RESTful API executing CRUD operations instantly over a high-fidelity HTML5 glassmorphism GUI.")
], Inches(1), Inches(1.8), Inches(11), Inches(5))


# --- SLIDE 5: Go To Market & Scalability ---
slide5 = prs.slides.add_slide(blank_layout)
apply_dark_theme(slide5)
title5 = slide5.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11), Inches(1))
format_title(title5, "AI-Readiness & Evolution")
add_body_text(slide5, [
    ("Enterprise SIEM Integration Ready", "Because CCAF utilizes standard REST endpoints (/api/block, /api/pcap), adapting to existing enterprise logs is completely frictionless."),
    ("Behavioral AI Pipeline", "Our structured packet routing naturally establishes a pipeline for next-gen integration with Machine Learning models (e.g., YOLOv8-based anomaly detection) to identify zero-day threats autonomously."),
    ("Target Deployment", "Ideal for critical infrastructure edge nodes, high-security local networks, and performance-sensitive proxy bounds.")
], Inches(1), Inches(1.8), Inches(11), Inches(5))

# Save the presentation
output_path = "CCAF_Investor_Pitch.pptx"
prs.save(output_path)
print(f"Investor Pitch Deck successfully generated at: {output_path}")
